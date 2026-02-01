#!/usr/bin/env python3
from __future__ import annotations

import base64
import re
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path("/Users/RoanV/Year-10-Digital-Tech")
RES_DIR = ROOT / "resources" / "term-1"
BRAND_DIR = ROOT / "resources" / "_brand"
LOGO_SRC = ROOT / "website" / "public" / "images" / "tkc-crest.svg"

TKC_BLUE = RGBColor(0x00, 0x38, 0x65)
TKC_GOLD = RGBColor(0xC5, 0xB7, 0x83)
TEXT_DARK = RGBColor(0x11, 0x18, 0x27)


def extract_logo_png() -> Path:
    BRAND_DIR.mkdir(parents=True, exist_ok=True)
    out_path = BRAND_DIR / "tkc-crest.png"
    if out_path.exists():
        return out_path
    svg_text = LOGO_SRC.read_text(encoding="utf-8", errors="ignore")
    match = re.search(r"data:image/png;base64,([A-Za-z0-9+/=]+)", svg_text)
    if not match:
        raise RuntimeError("Embedded PNG not found in tkc-crest.svg")
    png_bytes = base64.b64decode(match.group(1))
    out_path.write_bytes(png_bytes)
    return out_path


def add_brand_bar(slide, prs: Presentation, title: str) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TKC_BLUE
    bar.line.fill.background()

    text_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.05), prs.slide_width - Inches(0.6), Inches(0.3)
    )
    tf = text_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True


def add_title_slide(prs: Presentation, logo_path: Path, week: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs, "Year 10 Digital Technologies")

    slide.shapes.add_picture(str(logo_path), Inches(0.5), Inches(1.2), height=Inches(1.1))

    title_box = slide.shapes.add_textbox(
        Inches(1.9), Inches(1.2), prs.slide_width - Inches(2.4), Inches(1.2)
    )
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Term 1 – Week {week['num']:02d}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TKC_BLUE

    subtitle = slide.shapes.add_textbox(
        Inches(1.9), Inches(2.2), prs.slide_width - Inches(2.4), Inches(1.0)
    )
    st = subtitle.text_frame
    st.clear()
    p2 = st.paragraphs[0]
    p2.text = week["focus"]
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_DARK

    date_box = slide.shapes.add_textbox(
        Inches(1.9), Inches(2.9), prs.slide_width - Inches(2.4), Inches(0.6)
    )
    dt = date_box.text_frame
    dt.clear()
    p3 = dt.paragraphs[0]
    p3.text = week["dates"]
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_bar(slide, prs, "Year 10 Digital Technologies")

    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.7), prs.slide_width - Inches(1.2), Inches(0.6)
    )
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TKC_BLUE

    body_box = slide.shapes.add_textbox(
        Inches(0.9), Inches(1.6), prs.slide_width - Inches(1.5), prs.slide_height - Inches(2.0)
    )
    body = body_box.text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_DARK


def build_pptx(week: dict, out_path: Path, logo_path: Path) -> None:
    prs = Presentation()
    add_title_slide(prs, logo_path, week)
    add_bullet_slide(
        prs,
        "Lesson Objective & Overview",
        [
            f"Objective: {week['objective']}",
            f"Overview: {week['overview']}",
        ],
    )
    add_bullet_slide(prs, "Key Concepts", week["concepts"])
    add_bullet_slide(prs, "Lesson Sequence", week["sequence"])
    add_bullet_slide(prs, "Workshop / Activity", week["activity"])
    exit_bullets = [
        f"Evidence: {week['evidence'][0]}",
        f"Evidence: {week['evidence'][1]}",
    ]
    if week["checkpoint"]:
        exit_bullets.insert(0, f"Checkpoint: {week['checkpoint']}")
    exit_bullets.append(f"Outcome tags: {week['outcomes']}")
    exit_bullets.append(f"Exit ticket: {week['exit_ticket']}")
    add_bullet_slide(prs, "Checkpoint & Exit Ticket", exit_bullets)
    prs.save(out_path)


def worksheet_html(week: dict, logo_path: Path) -> str:
    logo_uri = f"file://{logo_path}"
    tasks = "".join(f"<li>{item}</li>" for item in week["worksheet"]["tasks"])
    extra = week["worksheet"]["extra_html"]
    return f"""<!DOCTYPE html>
<html>
<head>
  <meta charset="utf-8" />
  <style>
    @page {{ margin: 14mm; }}
    body {{ font-family: Arial, sans-serif; color: #111827; }}
    h1, h2 {{ color: #003865; }}
    .brand-header {{
      background: #003865;
      color: #fff;
      padding: 12px 14px;
      border-bottom: 4px solid #C5B783;
      margin: 0 0 12px;
    }}
    .brand-table {{ width: 100%; border-collapse: collapse; }}
    .brand-table td {{ border: none; }}
    .brand-logo {{ width: 48px; height: auto; }}
    .brand-title {{ font-weight: 700; font-size: 18px; }}
    .brand-subtitle {{ color: #C5B783; font-size: 12px; }}
    table {{ width: 100%; border-collapse: collapse; font-size: 12px; }}
    th, td {{ border: 1px solid #bbb; padding: 6px 8px; text-align: left; }}
    th {{ background: #f0f0f0; }}
    .box {{ border: 1px solid #999; height: 120px; margin: 8px 0; }}
    .small {{ font-size: 12px; color: #4b5563; }}
  </style>
</head>
<body>
  <div class="brand-header">
    <table class="brand-table">
      <tr>
        <td style="width:56px;"><img class="brand-logo" src="{logo_uri}" alt="The King's College" /></td>
        <td>
          <div class="brand-title">Year 10 Digital Technologies</div>
          <div class="brand-subtitle">Term 1 · Week {week['num']:02d}</div>
        </td>
      </tr>
    </table>
  </div>

  <h1>Week {week['num']:02d} Worksheet</h1>
  <p><strong>Focus:</strong> {week['focus']}</p>
  <p><strong>Objective:</strong> {week['objective']}</p>
  <p class="small"><strong>Outcome tags:</strong> {week['outcomes']}</p>

  <h2>Tasks</h2>
  <ul>{tasks}</ul>
  {extra}
</body>
</html>
"""


def resource_html(week: dict, logo_path: Path, title: str, body_html: str) -> str:
    logo_uri = f"file://{logo_path}"
    return f"""<!DOCTYPE html>
<html>
<head>
  <meta charset="utf-8" />
  <style>
    @page {{ margin: 14mm; }}
    body {{ font-family: Arial, sans-serif; color: #111827; }}
    h1, h2 {{ color: #003865; }}
    .brand-header {{
      background: #003865;
      color: #fff;
      padding: 12px 14px;
      border-bottom: 4px solid #C5B783;
      margin: 0 0 12px;
    }}
    .brand-table {{ width: 100%; border-collapse: collapse; }}
    .brand-table td {{ border: none; }}
    .brand-logo {{ width: 48px; height: auto; }}
    .brand-title {{ font-weight: 700; font-size: 18px; }}
    .brand-subtitle {{ color: #C5B783; font-size: 12px; }}
    table {{ width: 100%; border-collapse: collapse; font-size: 12px; }}
    th, td {{ border: 1px solid #bbb; padding: 6px 8px; text-align: left; }}
    th {{ background: #f0f0f0; }}
    .box {{ border: 1px solid #999; height: 120px; margin: 8px 0; }}
    .small {{ font-size: 12px; color: #4b5563; }}
  </style>
</head>
<body>
  <div class="brand-header">
    <table class="brand-table">
      <tr>
        <td style="width:56px;"><img class="brand-logo" src="{logo_uri}" alt="The King's College" /></td>
        <td>
          <div class="brand-title">Year 10 Digital Technologies</div>
          <div class="brand-subtitle">Term 1 · Week {week['num']:02d}</div>
        </td>
      </tr>
    </table>
  </div>

  <h1>{title}</h1>
  <p class="small"><strong>Focus:</strong> {week['focus']} | <strong>Outcome tags:</strong> {week['outcomes']}</p>
  {body_html}
</body>
</html>
"""


def write_worksheet(week: dict, out_dir: Path, logo_path: Path) -> None:
    html_path = out_dir / "worksheet.html"
    pdf_path = out_dir / "worksheet.pdf"
    html_path.write_text(worksheet_html(week, logo_path), encoding="utf-8")
    subprocess.run(["weasyprint", str(html_path), str(pdf_path)], check=True)

    extras = week.get("extras", [])
    for extra in extras:
        extra_html_path = out_dir / f"{extra['filename']}.html"
        extra_pdf_path = out_dir / f"{extra['filename']}.pdf"
        extra_html_path.write_text(
            resource_html(week, logo_path, extra["title"], extra["body_html"]),
            encoding="utf-8",
        )
        subprocess.run(["weasyprint", str(extra_html_path), str(extra_pdf_path)], check=True)


def build_resources() -> None:
    logo_path = extract_logo_png()
    for week in WEEKS:
        week_dir = RES_DIR / f"week-{week['num']:02d}"
        week_dir.mkdir(parents=True, exist_ok=True)
        pptx_path = week_dir / "slides.pptx"
        build_pptx(week, pptx_path, logo_path)
        write_worksheet(week, week_dir, logo_path)
        for extra_file in week.get("extra_files", []):
            (week_dir / extra_file["name"]).write_text(extra_file["content"], encoding="utf-8")


WEEKS = [
    {
        "num": 1,
        "dates": "2 Feb - 6 Feb",
        "focus": "Course launch, major project overview, digital systems intro",
        "checkpoint": "Baseline diagnostic",
        "outcomes": "DS10-1 (WA10DIGDS1)",
        "objective": "I can explain what this course is about, describe the major project, and show how data moves through a digital system.",
        "overview": "We will launch the course, unpack the major project, and complete a short digital systems diagnostic task.",
        "concepts": [
            "Course expectations and assessment structure",
            "Major project phases and milestones",
            "Digital system components (input, process, storage, output)",
            "Data flow and security points",
        ],
        "sequence": [
            "Map a secure data flow and identify potential attack points.",
            "Mini-lesson on key concepts and vocabulary.",
            "Guided practice or studio time aligned to the focus.",
        ],
        "activity": [
            "Choose a simple system (login, online order, ticket scan).",
            "Map data flow and identify at least 2 risks.",
            "Suggest one mitigation for each risk.",
        ],
        "evidence": [
            "Annotated data-flow diagram and threat notes.",
            "Classwork artefact or screenshot.",
        ],
        "exit_ticket": "One risk in a digital system and a mitigation.",
        "worksheet": {
            "tasks": [
                "Map a digital system data flow.",
                "Identify risks and propose mitigations.",
            ],
            "extra_html": (
                "<h2>Data Flow Map</h2>"
                "<table><tr><th>Component</th><th>Data In/Out</th><th>Risk</th><th>Mitigation</th></tr>"
                "<tr><td></td><td></td><td></td><td></td></tr>"
                "<tr><td></td><td></td><td></td><td></td></tr>"
                "<tr><td></td><td></td><td></td><td></td></tr>"
                "</table>"
            ),
        },
        "extras": [
            {
                "filename": "baseline-diagnostic",
                "title": "Baseline Diagnostic",
                "body_html": (
                    "<ol>"
                    "<li>Define input, process, storage, and output.</li>"
                    "<li>Give one example of a digital system.</li>"
                    "<li>Identify one risk in a data flow.</li>"
                    "<li>What is one way to mitigate that risk?</li>"
                    "<li>What do you want to learn in this course?</li>"
                    "</ol>"
                ),
            },
            {
                "filename": "major-project-overview",
                "title": "Major Project Overview",
                "body_html": (
                    "<h2>Goal</h2><p>Build a database‑driven website/app that demonstrates your skills in design, data, and implementation.</p>"
                    "<h2>Key deliverables</h2><ul>"
                    "<li>Design brief + success criteria</li>"
                    "<li>Prototype and evaluation</li>"
                    "<li>Working build with documentation</li>"
                    "</ul>"
                    "<h2>Milestones</h2><ul>"
                    "<li>Proposal and brief draft (Term 1)</li>"
                    "<li>Data model and early build (Term 2)</li>"
                    "<li>Core features and testing (Term 3)</li>"
                    "<li>Final delivery and reflection (Term 4)</li>"
                    "</ul>"
                ),
            },
        ],
    },
    {
        "num": 2,
        "dates": "9 Feb - 13 Feb",
        "focus": "Hardware specs, performance, suitability",
        "checkpoint": "",
        "outcomes": "DS10-2 (scope/sequence)",
        "objective": "I can compare hardware specifications and justify which device is best for a given task.",
        "overview": "We will compare CPU, RAM, storage, and GPU specs, then justify a device choice using evidence.",
        "concepts": [
            "CPU, RAM, storage, GPU, I/O basics",
            "Benchmarking and performance indicators",
            "Matching specs to task requirements",
        ],
        "sequence": [
            "Compare hardware specifications and justify device choices for a task.",
            "Mini-lesson on key concepts and vocabulary.",
            "Guided practice or studio time aligned to the focus.",
        ],
        "activity": [
            "Compare at least two devices using a table.",
            "Select the best device for a given task.",
            "Write a short justification using evidence.",
        ],
        "evidence": [
            "Hardware justification table.",
            "Classwork artefact or screenshot.",
        ],
        "exit_ticket": "Which component matters most for your chosen task and why?",
        "worksheet": {
            "tasks": [
                "Compare at least two devices using the table.",
                "Write a short justification for your choice.",
            ],
            "extra_html": (
                "<h2>Hardware Comparison</h2>"
                "<table><tr><th>Device</th><th>CPU</th><th>RAM</th><th>Storage</th><th>GPU</th><th>Price</th><th>Best Use</th></tr>"
                "<tr><td></td><td></td><td></td><td></td><td></td><td></td><td></td></tr>"
                "<tr><td></td><td></td><td></td><td></td><td></td><td></td><td></td></tr>"
                "</table>"
                "<h2>Justification</h2>"
                "<div class=\"box\"></div>"
            ),
        },
        "extras": [
            {
                "filename": "device-specs",
                "title": "Device Specification Sheet",
                "body_html": (
                    "<table><tr><th>Device</th><th>CPU</th><th>RAM</th><th>Storage</th><th>GPU</th><th>Notes</th></tr>"
                    "<tr><td>Device A</td><td></td><td></td><td></td><td></td><td></td></tr>"
                    "<tr><td>Device B</td><td></td><td></td><td></td><td></td><td></td></tr>"
                    "<tr><td>Device C</td><td></td><td></td><td></td><td></td><td></td></tr>"
                    "</table>"
                ),
            },
            {
                "filename": "benchmark-notes",
                "title": "Benchmark Notes",
                "body_html": (
                    "<ul>"
                    "<li>CPU: single‑core vs multi‑core performance</li>"
                    "<li>RAM: multitasking and file size limits</li>"
                    "<li>Storage: SSD vs HDD load times</li>"
                    "<li>GPU: rendering and graphics workloads</li>"
                    "</ul>"
                ),
            },
        ],
    },
    {
        "num": 3,
        "dates": "16 Feb - 20 Feb",
        "focus": "Web documents: content/structure/presentation (HTML/CSS)",
        "checkpoint": "",
        "outcomes": "DR10-1 (WA10DIGDR1)",
        "objective": "I can build a semantic HTML page and style it using a separate CSS file.",
        "overview": "We will learn the roles of HTML and CSS and build a simple page using semantic elements.",
        "concepts": [
            "HTML structure vs CSS presentation",
            "Semantic tags (header, nav, main, section, footer)",
            "CSS selectors and basic layout",
        ],
        "sequence": [
            "Build a semantic HTML page and separate CSS for structure vs presentation.",
            "Mini-lesson on key concepts and vocabulary.",
            "Guided practice or studio time aligned to the focus.",
        ],
        "activity": [
            "Sketch a simple page layout.",
            "Build the HTML structure using semantic tags.",
            "Apply at least five CSS styles.",
        ],
        "evidence": [
            "HTML/CSS page with explanation of structure/presentation.",
            "Classwork artefact or screenshot.",
        ],
        "exit_ticket": "List three semantic tags and their purpose.",
        "worksheet": {
            "tasks": [
                "Plan a page layout using semantic tags.",
                "List 5 CSS styles you will apply.",
            ],
            "extra_html": (
                "<h2>Page Plan</h2>"
                "<div class=\"box\"></div>"
                "<h2>Semantic Tags Checklist</h2>"
                "<ul><li>header</li><li>nav</li><li>main</li><li>section</li><li>footer</li></ul>"
            ),
        },
        "extras": [
            {
                "filename": "html-css-reference",
                "title": "HTML/CSS Reference",
                "body_html": (
                    "<h2>HTML</h2><ul>"
                    "<li>header, nav, main, section, article, footer</li>"
                    "<li>h1‑h6, p, ul/ol, a, img</li>"
                    "</ul>"
                    "<h2>CSS</h2><ul>"
                    "<li>Selectors: element, class, id</li>"
                    "<li>Box model: margin, border, padding</li>"
                    "<li>Layout: display, flex</li>"
                    "</ul>"
                ),
            },
        ],
        "extra_files": [
            {
                "name": "starter-template.html",
                "content": (
                    "<!doctype html>\\n"
                    "<html lang=\\\"en\\\">\\n"
                    "<head>\\n  <meta charset=\\\"utf-8\\\" />\\n"
                    "  <title>Starter Page</title>\\n"
                    "  <link rel=\\\"stylesheet\\\" href=\\\"styles.css\\\" />\\n"
                    "</head>\\n"
                    "<body>\\n"
                    "  <header><h1>Page Title</h1></header>\\n"
                    "  <nav>Navigation</nav>\\n"
                    "  <main>\\n    <section><h2>Section</h2><p>Content</p></section>\\n"
                    "  </main>\\n"
                    "  <footer>Footer</footer>\\n"
                    "</body>\\n"
                    "</html>\\n"
                ),
            },
            {
                "name": "styles.css",
                "content": (
                    "body { font-family: Arial, sans-serif; margin: 24px; }\\n"
                    "header, nav, main, footer { margin-bottom: 16px; }\\n"
                    "nav { background: #f3f4f6; padding: 8px; }\\n"
                    "section { border: 1px solid #e5e7eb; padding: 12px; }\\n"
                ),
            },
        ],
    },
    {
        "num": 4,
        "dates": "23 Feb - 27 Feb",
        "focus": "Client needs, data gathering for requirements",
        "checkpoint": "Minor A1 brief issued",
        "outcomes": "DI10-1 (WA10DIGDI1), ID10-1 (WA10DIGDTID1)",
        "objective": "I can gather client needs and define the problem, audience, and constraints.",
        "overview": "We will collect stakeholder data and draft a clear problem statement for A1.",
        "concepts": [
            "Stakeholders and client needs",
            "Research methods (survey, interview, observation)",
            "Problem framing and constraints",
        ],
        "sequence": [
            "Gather stakeholder data and define client needs and constraints.",
            "Define the problem and audience through evidence and interviews.",
            "Mini-lesson on key concepts and vocabulary.",
        ],
        "activity": [
            "Conduct a short interview or survey.",
            "Summarise key needs and constraints.",
            "Draft a problem statement.",
        ],
        "evidence": [
            "Client needs statement with evidence.",
            "Problem definition and audience summary.",
        ],
        "exit_ticket": "One client need and one constraint with evidence.",
        "worksheet": {
            "tasks": [
                "Complete the interview/survey prompts.",
                "Draft a problem statement and audience summary.",
            ],
            "extra_html": (
                "<h2>Interview / Survey Prompts</h2>"
                "<ol><li>What problem are we trying to solve?</li>"
                "<li>Who is the primary user?</li>"
                "<li>What tasks must the solution support?</li>"
                "<li>What constraints or limitations exist?</li></ol>"
                "<h2>Problem Statement</h2><div class=\"box\"></div>"
            ),
        },
        "extras": [
            {
                "filename": "survey-template",
                "title": "Survey / Interview Template",
                "body_html": (
                    "<ol>"
                    "<li>What problem are we trying to solve?</li>"
                    "<li>Who is the primary user?</li>"
                    "<li>What tasks must the solution support?</li>"
                    "<li>What constraints or limitations exist?</li>"
                    "</ol>"
                ),
            },
            {
                "filename": "problem-framing",
                "title": "Problem Framing Template",
                "body_html": (
                    "<p><strong>Problem statement:</strong></p><div class=\"box\"></div>"
                    "<p><strong>Target audience:</strong></p><div class=\"box\"></div>"
                    "<p><strong>Constraints:</strong></p><div class=\"box\"></div>"
                ),
            },
        ],
    },
    {
        "num": 5,
        "dates": "2 Mar - 6 Mar",
        "focus": "UX/UI design, wireframes, criteria",
        "checkpoint": "",
        "outcomes": "DI10-2 (WA10DIGDI2), DSN10-1 (WA10DIGDTDE1)",
        "objective": "I can design wireframes and evaluate options against criteria.",
        "overview": "We will build low-fidelity wireframes and test them against success criteria.",
        "concepts": [
            "UX vs UI and wireframes",
            "Design criteria and constraints",
            "Evaluating alternatives",
        ],
        "sequence": [
            "Design and prototype UX/UI screens based on user requirements.",
            "Sketch alternative solutions and evaluate them against criteria.",
            "Mini-lesson on key concepts and vocabulary.",
        ],
        "activity": [
            "Sketch two layout options.",
            "Build a wireframe for the best option.",
            "Evaluate using criteria.",
        ],
        "evidence": [
            "Prototype screens with feedback notes.",
            "Alternatives sketch set with criteria notes.",
        ],
        "exit_ticket": "One design decision and the criterion it meets.",
        "worksheet": {
            "tasks": [
                "Sketch two alternative layouts.",
                "Select the best option and explain why.",
            ],
            "extra_html": (
                "<h2>Wireframe Sketches</h2>"
                "<div class=\"box\"></div>"
                "<div class=\"box\"></div>"
                "<h2>Evaluation Matrix</h2>"
                "<table><tr><th>Option</th><th>Criteria Met</th><th>Notes</th></tr>"
                "<tr><td>A</td><td></td><td></td></tr>"
                "<tr><td>B</td><td></td><td></td></tr></table>"
            ),
        },
        "extras": [
            {
                "filename": "wireframe-grid",
                "title": "Wireframe Grid",
                "body_html": "<div class=\"box\"></div><div class=\"box\"></div>",
            },
            {
                "filename": "evaluation-matrix",
                "title": "Evaluation Matrix",
                "body_html": (
                    "<table><tr><th>Option</th><th>Criteria Met</th><th>Notes</th></tr>"
                    "<tr><td>A</td><td></td><td></td></tr>"
                    "<tr><td>B</td><td></td><td></td></tr></table>"
                ),
            },
        ],
    },
    {
        "num": 6,
        "dates": "9 Mar - 13 Mar",
        "focus": "Prototype build, iterate on feedback",
        "checkpoint": "",
        "outcomes": "DI10-2 (WA10DIGDI2), EV10-1 (WA10DIGDTEV1)",
        "objective": "I can improve a prototype using feedback and evaluation criteria.",
        "overview": "We will gather feedback, record improvements, and update prototypes.",
        "concepts": [
            "Feedback collection methods",
            "Evaluation criteria",
            "Iteration and change logs",
        ],
        "sequence": [
            "Design and prototype UX/UI screens based on user requirements.",
            "Evaluate the solution against criteria and record improvements.",
            "Mini-lesson on key concepts and vocabulary.",
        ],
        "activity": [
            "Collect feedback from peers.",
            "Record changes and update prototype.",
            "Summarise improvements.",
        ],
        "evidence": [
            "Prototype screens with feedback notes.",
            "Evaluation notes and improvement list.",
        ],
        "exit_ticket": "Top improvement you made and why.",
        "worksheet": {
            "tasks": [
                "Collect feedback using the form.",
                "Record at least 3 changes in the iteration log.",
            ],
            "extra_html": (
                "<h2>Feedback Form</h2>"
                "<ul><li>What worked well?</li><li>What was confusing?</li><li>One improvement to make.</li></ul>"
                "<h2>Iteration Log</h2>"
                "<table><tr><th>Change</th><th>Reason</th><th>Result</th></tr>"
                "<tr><td></td><td></td><td></td></tr>"
                "<tr><td></td><td></td><td></td></tr>"
                "<tr><td></td><td></td><td></td></tr></table>"
            ),
        },
        "extras": [
            {
                "filename": "feedback-form",
                "title": "Feedback Form",
                "body_html": (
                    "<ul><li>What worked well?</li><li>What was confusing?</li><li>One improvement to make.</li></ul>"
                ),
            },
            {
                "filename": "iteration-log",
                "title": "Iteration Log",
                "body_html": (
                    "<table><tr><th>Change</th><th>Reason</th><th>Result</th></tr>"
                    "<tr><td></td><td></td><td></td></tr>"
                    "<tr><td></td><td></td><td></td></tr></table>"
                ),
            },
        ],
    },
    {
        "num": 7,
        "dates": "16 Mar - 20 Mar",
        "focus": "Design brief drafting, project planning",
        "checkpoint": "",
        "outcomes": "ID10-2 (WA10DIGDTID2), PM10-1 (WA10DIGDTPM1)",
        "objective": "I can draft a design brief and plan tasks using a sprint board.",
        "overview": "We will write a clear design brief and plan tasks, risks, and time estimates.",
        "concepts": [
            "Design brief sections",
            "Success criteria and constraints",
            "Planning tools (Kanban, sprint board)",
        ],
        "sequence": [
            "Draft a design brief with success criteria and constraints.",
            "Plan work in a sprint board with tasks, risks and time estimates.",
            "Mini-lesson on key concepts and vocabulary.",
        ],
        "activity": [
            "Draft the design brief.",
            "Create a sprint plan with tasks and estimates.",
            "Identify one risk and mitigation.",
        ],
        "evidence": [
            "Design brief document.",
            "Updated project plan and task board.",
        ],
        "exit_ticket": "Three tasks you will complete next week.",
        "worksheet": {
            "tasks": [
                "Complete the design brief template.",
                "Create a sprint plan table.",
            ],
            "extra_html": (
                "<h2>Design Brief Template</h2>"
                "<ul><li>Problem statement</li><li>Target audience</li><li>Success criteria</li><li>Constraints</li><li>Proposed solution</li></ul>"
                "<h2>Sprint Plan</h2>"
                "<table><tr><th>Task</th><th>Owner</th><th>Estimate</th><th>Risk</th></tr>"
                "<tr><td></td><td></td><td></td><td></td></tr>"
                "<tr><td></td><td></td><td></td><td></td></tr></table>"
            ),
        },
        "extras": [
            {
                "filename": "design-brief-template",
                "title": "Design Brief Template",
                "body_html": (
                    "<ul><li>Problem statement</li><li>Target audience</li><li>Success criteria</li>"
                    "<li>Constraints</li><li>Proposed solution</li></ul>"
                    "<div class=\"box\"></div>"
                ),
            },
            {
                "filename": "sprint-plan",
                "title": "Sprint Plan",
                "body_html": (
                    "<table><tr><th>Task</th><th>Owner</th><th>Estimate</th><th>Risk</th></tr>"
                    "<tr><td></td><td></td><td></td><td></td></tr>"
                    "<tr><td></td><td></td><td></td><td></td></tr></table>"
                ),
            },
        ],
    },
    {
        "num": 8,
        "dates": "23 Mar - 27 Mar",
        "focus": "Present prototype and rationale",
        "checkpoint": "Major project proposal kickoff",
        "outcomes": "EV10-1 (WA10DIGDTEV1), PM10-1 (WA10DIGDTPM1)",
        "objective": "I can present my prototype, explain design decisions, and update my plan.",
        "overview": "We will present prototypes, evaluate against criteria, and refine plans.",
        "concepts": [
            "Presentation structure",
            "Evidence-based evaluation",
            "Next steps and planning updates",
        ],
        "sequence": [
            "Evaluate the solution against criteria and record improvements.",
            "Plan work in a sprint board with tasks, risks and time estimates.",
            "Mini-lesson on key concepts and vocabulary.",
        ],
        "activity": [
            "Present prototype and rationale.",
            "Collect peer feedback.",
            "Update plan based on feedback.",
        ],
        "evidence": [
            "Evaluation notes and improvement list.",
            "Updated project plan and task board.",
        ],
        "exit_ticket": "Best feedback received and next action.",
        "worksheet": {
            "tasks": [
                "Use the checklist to prepare your presentation.",
                "Complete the evaluation summary table.",
            ],
            "extra_html": (
                "<h2>Presentation Checklist</h2>"
                "<ul><li>Problem and audience</li><li>Prototype walkthrough</li><li>Design decisions</li><li>Next steps</li></ul>"
                "<h2>Evaluation Summary</h2>"
                "<table><tr><th>Criteria</th><th>Evidence</th><th>Improvement</th></tr>"
                "<tr><td></td><td></td><td></td></tr>"
                "<tr><td></td><td></td><td></td></tr></table>"
            ),
        },
        "extras": [
            {
                "filename": "presentation-checklist",
                "title": "Presentation Checklist",
                "body_html": (
                    "<ul><li>Problem and audience</li><li>Prototype walkthrough</li>"
                    "<li>Design decisions</li><li>Next steps</li></ul>"
                ),
            },
            {
                "filename": "evaluation-summary",
                "title": "Evaluation Summary",
                "body_html": (
                    "<table><tr><th>Criteria</th><th>Evidence</th><th>Improvement</th></tr>"
                    "<tr><td></td><td></td><td></td></tr>"
                    "<tr><td></td><td></td><td></td></tr></table>"
                ),
            },
        ],
    },
    {
        "num": 9,
        "dates": "30 Mar - 2 Apr",
        "focus": "Consolidation and reflections",
        "checkpoint": "Minor A1 due; Major project brief draft",
        "outcomes": "ID10-2 (WA10DIGDTID2), ID10-3 (WA10DIGDTID3)",
        "objective": "I can finalise my design brief and justify technology choices.",
        "overview": "We will submit A1 and refine the major project brief with evidence and trade-offs.",
        "concepts": [
            "Finalising a design brief",
            "Technology trade-offs",
            "Reflection and goal setting",
        ],
        "sequence": [
            "Draft a design brief with success criteria and constraints.",
            "Investigate technologies and constraints and summarise trade-offs.",
            "Mini-lesson on key concepts and vocabulary.",
        ],
        "activity": [
            "Finalise design brief and submit A1.",
            "Complete a trade-off table.",
            "Write a reflection on learning.",
        ],
        "evidence": [
            "Design brief document.",
            "Technology comparison summary.",
        ],
        "exit_ticket": "One trade-off and why you chose it.",
        "worksheet": {
            "tasks": [
                "Complete the final checklist.",
                "Fill in the trade-off table.",
                "Write a short reflection.",
            ],
            "extra_html": (
                "<h2>Final Checklist</h2>"
                "<ul><li>Problem statement clear</li><li>Audience defined</li><li>Success criteria listed</li><li>Constraints acknowledged</li><li>Proposed solution described</li></ul>"
                "<h2>Technology Trade-off Table</h2>"
                "<table><tr><th>Option</th><th>Benefit</th><th>Risk</th><th>Decision</th></tr>"
                "<tr><td></td><td></td><td></td><td></td></tr>"
                "<tr><td></td><td></td><td></td><td></td></tr></table>"
                "<h2>Reflection</h2><div class=\"box\"></div>"
            ),
        },
        "extras": [
            {
                "filename": "final-checklist",
                "title": "Final Checklist",
                "body_html": (
                    "<ul><li>Problem statement clear</li><li>Audience defined</li>"
                    "<li>Success criteria listed</li><li>Constraints acknowledged</li>"
                    "<li>Proposed solution described</li></ul>"
                ),
            },
            {
                "filename": "tradeoff-table",
                "title": "Technology Trade-off Table",
                "body_html": (
                    "<table><tr><th>Option</th><th>Benefit</th><th>Risk</th><th>Decision</th></tr>"
                    "<tr><td></td><td></td><td></td><td></td></tr>"
                    "<tr><td></td><td></td><td></td><td></td></tr></table>"
                ),
            },
            {
                "filename": "reflection-prompt",
                "title": "Reflection Prompt",
                "body_html": (
                    "<p>What did you learn this term? What will you improve next term?</p>"
                    "<div class=\"box\"></div>"
                ),
            },
        ],
    },
]


if __name__ == "__main__":
    build_resources()
